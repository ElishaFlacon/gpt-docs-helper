import os
import pptx
import docx
from PyPDF2 import PdfReader
from langchain.llms import OpenAI
from langchain.vectorstores import FAISS
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.chains.question_answering import load_qa_chain


class AI:
    def __init__(self):
        self.directory_path = os.getenv(
            "DOCUMENTS_DIRECTORY_PATH",
            "documents",
        )

        self.embeddings = OpenAIEmbeddings()
        self.llm = OpenAI(temperature=0)
        self.chain = load_qa_chain(self.llm, chain_type="stuff")

        self.docsearch = None
        self.count_documents = 0

    def upload_file(self) -> any:
        if os.path.exists(self.directory_path):
            file_paths = [
                os.path.join(self.directory_path, f)
                for f in os.listdir(self.directory_path)
                if os.path.isfile(os.path.join(self.directory_path, f))
            ]
            return file_paths
        else:
            print("Directory not found!")
            return []

    def extract_texts(self, root_files: any) -> FAISS:
        raw_text = ''

        for root_file in root_files:
            _, ext = os.path.splitext(root_file)
            if ext == '.pdf':
                with open(root_file, 'rb') as f:
                    reader = PdfReader(f)
                    for i in range(len(reader.pages)):
                        page = reader.pages[i]
                        raw_text += page.extract_text()
            elif ext == '.docx':
                doc = docx.Document(root_file)
                for paragraph in doc.paragraphs:
                    raw_text += paragraph.text
            elif ext == '.pptx':
                ppt = pptx.Presentation(root_file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            raw_text += shape.text

        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

        texts = text_splitter.split_text(raw_text)

        docsearch = FAISS.from_texts(texts, self.embeddings)
        return docsearch

    def extract_documents(self) -> None:
        root_files = self.upload_file()
        if not root_files:
            print("No files uploaded. Exiting.")
            return

        self.count_documents = len(root_files)
        self.docsearch = self.extract_texts(root_files)

    def query(self, query: str) -> str:
        docs = self.docsearch.similarity_search(query)

        russian_query = f"Дай ответ на этот вопрос на русском языке: {query}"

        return self.chain.run(
            input_documents=docs,
            question=russian_query,
        )
